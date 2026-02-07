"""
PowerPoint document processor.
Extracts text content from PowerPoint presentations.
"""
import io
from typing import Optional
from pptx import Presentation
from .base import BaseProcessor


class PowerPointProcessor(BaseProcessor):
    """Processor for PowerPoint documents (PPTX, PPT)"""
    
    def extract_text(self, file_content: bytes, filename: str) -> Optional[str]:
        """
        Extract text from PowerPoint file.
        
        Args:
            file_content: PowerPoint file content as bytes
            filename: Original filename
            
        Returns:
            Extracted text content, or None if extraction fails
        """
        try:
            ppt_file = io.BytesIO(file_content)
            prs = Presentation(ppt_file)
            
            text_parts = []
            
            # Extract text from each slide
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_text.append(f"[Slide {slide_num}]")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Extract text from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(" | ".join(row_text))
                
                if len(slide_text) > 1:  # More than just the slide number
                    text_parts.append("\n".join(slide_text))
            
            if not text_parts:
                return None
            
            full_text = "\n\n".join(text_parts)
            return self.clean_text(full_text)
        
        except Exception as e:
            print(f"Error processing PowerPoint {filename}: {e}")
            return None
